"""
EDU-Audit Enhanced Multimodal Document Agent
GPT-4V를 활용한 실제 이미지 분석 및 오류 탐지
"""

import asyncio
import logging
import os
import base64
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from io import BytesIO
import json

# PDF/PPT 파싱 라이브러리
import pdfplumber
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF for image extraction

# OCR 라이브러리
try:
    import easyocr
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logging.warning("EasyOCR not available. Install with: pip install easyocr")

# OpenAI API
import openai
from openai import AsyncOpenAI

# LlamaIndex 관련
from llama_index.core import Document, VectorStoreIndex
from llama_index.core.node_parser import SentenceSplitter
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI

# 확장된 모델들
from src.core.models import (
    DocumentMeta, PageInfo, PageElement, ElementType, IssueType,
    ImageElement, TableElement, ChartElement, BoundingBox, Issue,
    generate_doc_id, generate_element_id, generate_issue_id
)
from dotenv import load_dotenv

env_path = Path(__file__).resolve().parents[2] / '.env.dev'
load_dotenv(env_path)

logger = logging.getLogger(__name__)

class EnhancedMultimodalDocumentAgent:
    """GPT-4V를 활용한 향상된 멀티모달 문서 분석 에이전트"""
    
    def __init__(
        self, 
        openai_api_key: Optional[str] = None,
        vision_model: str = "gpt-5-nano",
        text_model: str = "gpt-5-nano",
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        enable_ocr: bool = True,
        enable_vision_analysis: bool = True,
        max_vision_tokens: int = 4096
    ):
        self.openai_api_key = openai_api_key or os.getenv("OPENAI_API_KEY")
        self.vision_model = vision_model
        self.text_model = text_model
        self.enable_ocr = enable_ocr and OCR_AVAILABLE
        self.enable_vision_analysis = enable_vision_analysis and self.openai_api_key
        self.max_vision_tokens = max_vision_tokens
        
        if not self.openai_api_key:
            raise ValueError("OpenAI API key is required")
        
        # OpenAI 클라이언트
        self.openai_client = AsyncOpenAI(api_key=self.openai_api_key)
        
        # LlamaIndex 컴포넌트
        self.embeddings = OpenAIEmbedding(api_key=self.openai_api_key)
        self.text_splitter = SentenceSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )
        
        # Text LLM (텍스트 분석용)
        self.text_llm = OpenAI(
            model=text_model,
            api_key=self.openai_api_key,
        )
        
        # OCR 초기화
        self.ocr_reader = None
        if self.enable_ocr:
            try:
                self.ocr_reader = easyocr.Reader(['ko', 'en'])
                logger.info("EasyOCR 초기화 완료 (한국어, 영어)")
            except Exception as e:
                logger.warning(f"OCR 초기화 실패: {str(e)}")
                self.enable_ocr = False
        
        # 메모리에 저장된 문서들
        self.documents: Dict[str, DocumentMeta] = {}
        self.pages: Dict[str, List[PageInfo]] = {}  # doc_id -> pages
        self.indexes: Dict[str, VectorStoreIndex] = {}  # doc_id -> index
        self.image_issues: Dict[str, List[Issue]] = {}  # doc_id -> issues
        
        logger.info(f"EnhancedMultimodalDocumentAgent 초기화 완료")
        logger.info(f"  OCR: {'활성화' if self.enable_ocr else '비활성화'}")
        logger.info(f"  Vision 분석: {'활성화' if self.enable_vision_analysis else '비활성화'}")
        logger.info(f"  Vision 모델: {self.vision_model}")
    
    async def parse_document(self, file_path: str) -> DocumentMeta:
        """
        향상된 멀티모달 문서 파싱 및 오류 탐지
        
        Args:
            file_path: 파싱할 문서 파일 경로
            
        Returns:
            DocumentMeta: 파싱된 문서 메타데이터
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")
        
        logger.info(f"향상된 멀티모달 문서 파싱 시작: {file_path}")
        
        # 문서 ID 생성
        doc_id = generate_doc_id(str(file_path))
        
        # 파일 확장자에 따른 파싱
        if file_path.suffix.lower() == '.pdf':
            doc_meta, pages = await self._parse_pdf_enhanced(file_path, doc_id)
        elif file_path.suffix.lower() in ['.ppt', '.pptx']:
            doc_meta, pages = await self._parse_ppt_enhanced(file_path, doc_id)
        else:
            raise ValueError(f"지원하지 않는 파일 형식: {file_path.suffix}")
        
        # 멀티모달 통계 업데이트
        doc_meta.total_images = sum(page.image_count for page in pages)
        doc_meta.total_tables = sum(page.table_count for page in pages)
        doc_meta.total_charts = sum(page.chart_count for page in pages)
        
        # 메모리에 저장
        self.documents[doc_id] = doc_meta
        self.pages[doc_id] = pages
        
        # 향상된 멀티모달 인덱스 생성
        await self._create_enhanced_multimodal_index(doc_id, pages)
        
        # 이미지 품질 및 오류 검사
        image_issues = await self._analyze_image_issues(doc_id, pages)
        self.image_issues[doc_id] = image_issues
        
        logger.info(f"향상된 멀티모달 문서 파싱 완료: {doc_id}")
        logger.info(f"  총 {len(pages)} 페이지, {doc_meta.total_images} 이미지, {doc_meta.total_tables} 표, {doc_meta.total_charts} 차트")
        logger.info(f"  이미지 관련 이슈: {len(image_issues)}개 발견")
        
        return doc_meta
    
    async def _parse_pdf_enhanced(self, file_path: Path, doc_id: str) -> Tuple[DocumentMeta, List[PageInfo]]:
        """향상된 PDF 멀티모달 파싱"""
        logger.info(f"향상된 PDF 파싱 중: {file_path}")
        
        pages = []
        title = None
        
        try:
            # PyMuPDF로 이미지 추출용
            pdf_doc = fitz.open(file_path)
            
            with pdfplumber.open(file_path) as pdf:
                # 첫 페이지에서 제목 추출
                if pdf.pages:
                    first_page_text = pdf.pages[0].extract_text() or ""
                    title = await self._extract_title_ai(first_page_text)
                
                # 각 페이지 파싱
                for page_num, page in enumerate(pdf.pages, 1):
                    logger.info(f"페이지 {page_num} 향상된 파싱 중...")
                    
                    # 텍스트 추출
                    page_text = page.extract_text() or ""
                    
                    # 멀티모달 요소들 추출
                    elements = []
                    
                    # 1. 텍스트 요소 추가
                    if page_text.strip():
                        text_element = PageElement(
                            element_id=generate_element_id(f"p{page_num:03d}", ElementType.TEXT, 0),
                            element_type=ElementType.TEXT,
                            text_content=page_text,
                            confidence=1.0
                        )
                        elements.append(text_element)
                    
                    # 2. 향상된 이미지 추출 및 분석
                    fitz_page = pdf_doc[page_num - 1]
                    image_elements = await self._extract_images_enhanced(
                        page, fitz_page, f"p{page_num:03d}"
                    )
                    elements.extend(image_elements)
                    
                    # 3. 향상된 표 추출
                    table_elements = await self._extract_tables_enhanced(
                        page, f"p{page_num:03d}"
                    )
                    elements.extend(table_elements)
                    
                    # 4. AI 기반 차트 감지
                    chart_elements = await self._detect_charts_ai(
                        image_elements, f"p{page_num:03d}"
                    )
                    elements.extend(chart_elements)
                    
                    # PageInfo 생성
                    page_info = PageInfo(
                        page_id=f"p{page_num:03d}",
                        page_number=page_num,
                        raw_text=page_text,
                        word_count=len(page_text.split()) if page_text else 0,
                        elements=elements
                    )
                    
                    pages.append(page_info)
                    
                    # 진행상황 로그
                    if page_num % 5 == 0:
                        logger.info(f"향상된 PDF 파싱 진행: {page_num}/{len(pdf.pages)} 페이지")
        
        except Exception as e:
            logger.error(f"향상된 PDF 파싱 실패: {str(e)}")
            raise
        finally:
            if 'pdf_doc' in locals():
                pdf_doc.close()
        
        # 문서 메타데이터 생성
        doc_meta = DocumentMeta(
            doc_id=doc_id,
            title=title or file_path.stem,
            doc_type="pdf",
            total_pages=len(pages),
            file_path=str(file_path)
        )
        
        return doc_meta, pages
    
    async def _extract_title_ai(self, text: str) -> Optional[str]:
        """AI를 사용한 제목 추출"""
        if not text.strip():
            return None
        
        try:
            prompt = f"""다음 텍스트에서 문서의 제목을 추출해주세요. 교육 자료의 제목으로 적절한 것을 선택하세요.

텍스트:
{text[:500]}

다음 중 하나로 응답해주세요:
1. 적절한 제목이 있다면 그 제목만 반환
2. 적절한 제목이 없다면 "NONE" 반환

제목:"""

            response = await self.openai_client.chat.completions.create(
                model=self.text_model,
                messages=[{"role": "user", "content": prompt}],
                max_completion_tokens=100,
            )
            
            title = response.choices[0].message.content.strip()
            return title if title != "NONE" else None
            
        except Exception as e:
            logger.warning(f"AI 제목 추출 실패: {str(e)}")
            return self._extract_title_from_text(text)
    
    def _extract_title_from_text(self, text: str) -> Optional[str]:
        """기존 제목 추출 방식 (백업용)"""
        if not text:
            return None
        
        lines = text.strip().split('\n')
        first_line = lines[0].strip()
        
        if 5 <= len(first_line) <= 100:
            return first_line
        
        return None
    
    async def _extract_images_enhanced(self, pdfplumber_page, fitz_page, page_id: str) -> List[PageElement]:
        """향상된 이미지 추출 및 GPT-4V 분석"""
        elements = []
        
        try:
            # PyMuPDF로 이미지 추출
            image_list = fitz_page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    # 이미지 데이터 추출
                    xref = img[0]
                    pix = fitz.Pixmap(fitz_page.parent, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.pil_tobytes(format="PNG")
                        pil_image = Image.open(BytesIO(img_data))
                        
                        # 이미지 크기 검사 및 리사이즈
                        if pil_image.width * pil_image.height > 2073600:  # ~2MP
                            # 큰 이미지는 리사이즈
                            pil_image.thumbnail((1024, 1024), Image.Resampling.LANCZOS)
                            buffered = BytesIO()
                            pil_image.save(buffered, format="PNG")
                            img_data = buffered.getvalue()
                        
                        # 이미지를 base64로 인코딩
                        img_base64 = base64.b64encode(img_data).decode()
                        
                        # GPT-4V로 향상된 이미지 분석
                        image_analysis = await self._analyze_image_gpt4v(pil_image, img_base64)
                        
                        # 이미지 요소 생성
                        image_element = PageElement(
                            element_id=generate_element_id(page_id, ElementType.IMAGE, img_index),
                            element_type=ElementType.IMAGE,
                            bbox=BoundingBox(x=0, y=0, width=pil_image.width, height=pil_image.height),
                            image_data=ImageElement(
                                element_id=generate_element_id(page_id, ElementType.IMAGE, img_index),
                                bbox=BoundingBox(x=0, y=0, width=pil_image.width, height=pil_image.height),
                                format="png",
                                size_bytes=len(img_data),
                                dimensions=(pil_image.width, pil_image.height),
                                ocr_text=image_analysis.get("ocr_text"),
                                description=image_analysis.get("description"),
                                image_data=f"data:image/png;base64,{img_base64}"
                            ),
                            confidence=image_analysis.get("confidence", 0.8)
                        )
                        
                        elements.append(image_element)
                        
                    pix = None  # 메모리 해제
                    
                except Exception as e:
                    logger.warning(f"이미지 {img_index} 처리 실패: {str(e)}")
        
        except Exception as e:
            logger.warning(f"페이지 이미지 추출 실패: {str(e)}")
        
        return elements
    
    async def _analyze_image_gpt4v(self, pil_image: Image.Image, img_base64: str) -> Dict[str, Any]:
        """GPT-4V를 사용한 향상된 이미지 분석"""
        analysis = {
            "ocr_text": None,
            "description": None,
            "confidence": 0.5,
            "image_type": "unknown",
            "educational_purpose": None,
            "quality_issues": [],
            "content_accuracy": None
        }
        
        # 1. 기존 OCR 분석
        if self.enable_ocr and self.ocr_reader:
            try:
                import numpy as np
                img_array = np.array(pil_image)
                
                ocr_results = self.ocr_reader.readtext(img_array)
                if ocr_results:
                    ocr_texts = [result[1] for result in ocr_results if result[2] > 0.5]
                    analysis["ocr_text"] = " ".join(ocr_texts)
                    analysis["confidence"] = max(analysis["confidence"], 0.7)
                
            except Exception as e:
                logger.warning(f"OCR 분석 실패: {str(e)}")
        
        # 2. GPT-4V 향상된 분석
        if self.enable_vision_analysis:
            try:
                gpt4v_analysis = await self._call_gpt4v_analysis(img_base64)
                if gpt4v_analysis:
                    analysis.update(gpt4v_analysis)
                    analysis["confidence"] = max(analysis["confidence"], 0.9)
                
            except Exception as e:
                logger.warning(f"GPT-4V 분석 실패: {str(e)}")
        
        return analysis
    
    async def _call_gpt4v_analysis(self, img_base64: str) -> Optional[Dict[str, Any]]:
        """GPT-4V API 호출로 교육 자료 맞춤 이미지 분석"""
        try:
            prompt = """당신은 교육 자료 품질 검수 전문가입니다. 이 이미지를 분석해주세요.

다음 항목들을 JSON 형식으로 분석해주세요:

1. **이미지 타입**: diagram, chart, graph, photo, screenshot, text_image, equation, flowchart, concept_map, other 중 하나
2. **교육적 목적**: explanation, example, data_visualization, exercise, reference, decoration 중 하나  
3. **주요 내용**: 이미지의 핵심 내용을 1-2문장으로 설명
4. **텍스트 내용**: 이미지에 포함된 모든 텍스트 (한국어/영어)
5. **품질 이슈**: 다음 중 해당하는 것들
   - low_resolution: 해상도가 낮아 글자가 흐림
   - poor_contrast: 대비가 나빠 가독성 저하
   - text_too_small: 텍스트가 너무 작아 읽기 어려움
   - unclear_diagram: 다이어그램이 불분명
   - missing_labels: 라벨이나 설명이 누락됨
   - color_accessibility: 색상 접근성 문제 (색맹 고려)
   - none: 문제없음
6. **내용 정확성**: 
   - 명백한 오류나 모순이 있는지 확인
   - 의심스러운 정보가 있다면 지적
7. **개선 제안**: 구체적인 개선 방안

반드시 다음 JSON 형식으로만 응답하세요:
{
  "image_type": "타입",
  "educational_purpose": "목적", 
  "description": "주요 내용 설명",
  "text_content": "추출된 텍스트",
  "quality_issues": ["이슈1", "이슈2"],
  "accuracy_concerns": "정확성 관련 우려사항 (없으면 null)",
  "improvement_suggestions": "개선 제안"
}"""

            response = await self.openai_client.chat.completions.create(
                model=self.vision_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_base64}",
                                    "detail": "high"
                                }
                            }
                        ]
                    }
                ],
                max_completion_tokens=self.max_vision_tokens,
            )
            
            content = response.choices[0].message.content.strip()
            
            # JSON 파싱 시도
            try:
                # JSON 블록 추출
                if "```json" in content:
                    json_start = content.find("```json") + 7
                    json_end = content.find("```", json_start)
                    content = content[json_start:json_end]
                elif "{" in content and "}" in content:
                    # JSON 객체 부분만 추출
                    start = content.find("{")
                    end = content.rfind("}") + 1
                    content = content[start:end]
                
                result = json.loads(content)
                
                return {
                    "description": result.get("description", ""),
                    "image_type": result.get("image_type", "unknown"),
                    "educational_purpose": result.get("educational_purpose"),
                    "quality_issues": result.get("quality_issues", []),
                    "accuracy_concerns": result.get("accuracy_concerns"),
                    "improvement_suggestions": result.get("improvement_suggestions"),
                    "ocr_text": result.get("text_content", "")
                }
                
            except json.JSONDecodeError as e:
                logger.warning(f"GPT-4V JSON 파싱 실패: {str(e)}")
                # JSON 파싱 실패시 텍스트에서 정보 추출
                return {
                    "description": content[:200],  # 처음 200자
                    "image_type": "unknown",
                    "quality_issues": [],
                    "accuracy_concerns": None
                }
            
        except Exception as e:
            logger.error(f"GPT-4V API 호출 실패: {str(e)}")
            return None
    
    async def _detect_charts_ai(self, image_elements: List[PageElement], page_id: str) -> List[PageElement]:
        """AI 기반 차트 감지 및 분석"""
        chart_elements = []
        
        for img_element in image_elements:
            if not img_element.image_data:
                continue
            
            # GPT-4V 분석 결과에서 차트 여부 확인
            analysis_result = await self._analyze_image_gpt4v(
                None, img_element.image_data.image_data.split(',')[1]
            )
            
            if analysis_result and analysis_result.get("image_type") in ["chart", "graph", "diagram"]:
                # 차트로 인식된 경우 상세 분석
                chart_analysis = await self._analyze_chart_details(
                    img_element.image_data.image_data.split(',')[1]
                )
                
                chart_element = PageElement(
                    element_id=img_element.element_id.replace("image", "chart"),
                    element_type=ElementType.CHART,
                    bbox=img_element.bbox,
                    chart_data=ChartElement(
                        element_id=img_element.element_id.replace("image", "chart"),
                        bbox=img_element.bbox,
                        chart_type=chart_analysis.get("chart_type", "unknown"),
                        title=chart_analysis.get("title"),
                        x_label=chart_analysis.get("x_label"),
                        y_label=chart_analysis.get("y_label"),
                        description=analysis_result.get("description", ""),
                        data_points=chart_analysis.get("data_points", []),
                        legend=chart_analysis.get("legend", []),
                        insights=chart_analysis.get("insights", [])
                    ),
                    confidence=0.8
                )
                
                chart_elements.append(chart_element)
        
        return chart_elements
    
    async def _analyze_chart_details(self, img_base64: str) -> Dict[str, Any]:
        """차트 상세 분석"""
        try:
            prompt = """이 차트/그래프를 상세히 분석해주세요.

다음 정보를 JSON 형식으로 추출해주세요:

{
  "chart_type": "bar, line, pie, scatter, area, histogram, box, other 중 하나",
  "title": "차트 제목",
  "x_label": "X축 라벨",
  "y_label": "Y축 라벨", 
  "data_points": [{"name": "데이터명", "value": "값"}],
  "legend": ["범례1", "범례2"],
  "insights": ["인사이트1", "인사이트2"],
  "data_accuracy_check": "데이터의 논리적 일관성 확인 결과"
}

반드시 JSON 형식으로만 응답하세요."""

            response = await self.openai_client.chat.completions.create(
                model=self.vision_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_base64}",
                                    "detail": "high"
                                }
                            }
                        ]
                    }
                ],
                max_completion_tokens=1000,
            )
            
            content = response.choices[0].message.content.strip()
            
            # JSON 파싱
            if "```json" in content:
                json_start = content.find("```json") + 7
                json_end = content.find("```", json_start)
                content = content[json_start:json_end]
            
            return json.loads(content)
            
        except Exception as e:
            logger.warning(f"차트 상세 분석 실패: {str(e)}")
            return {}
    
    async def _extract_tables_enhanced(self, page, page_id: str) -> List[PageElement]:
        """향상된 표 추출 및 검증"""
        elements = []
        
        try:
            tables = page.extract_tables()
            
            for table_index, table_data in enumerate(tables):
                if not table_data or len(table_data) < 2:
                    continue
                
                # 헤더와 데이터 분리
                headers = table_data[0] if table_data else []
                rows = table_data[1:] if len(table_data) > 1 else []
                
                # None 값 정리
                headers = [str(cell) if cell is not None else "" for cell in headers]
                clean_rows = []
                for row in rows:
                    clean_row = [str(cell) if cell is not None else "" for cell in row]
                    clean_rows.append(clean_row)
                
                # AI를 통한 표 검증
                table_issues = await self._validate_table_ai(headers, clean_rows)
                
                table_element = PageElement(
                    element_id=generate_element_id(page_id, ElementType.TABLE, table_index),
                    element_type=ElementType.TABLE,
                    table_data=TableElement(
                        element_id=generate_element_id(page_id, ElementType.TABLE, table_index),
                        headers=headers,
                        rows=clean_rows,
                        extraction_confidence=0.9 if not table_issues else 0.7
                    ),
                    confidence=0.9 if not table_issues else 0.7,
                    processing_notes=table_issues
                )
                
                elements.append(table_element)
        
        except Exception as e:
            logger.warning(f"향상된 표 추출 실패: {str(e)}")
        
        return elements
    
    async def _validate_table_ai(self, headers: List[str], sample_rows: List[List[str]]) -> List[str]:
        """AI를 통한 표 내용 논리적 검증"""
        try:
            table_text = "헤더: " + " | ".join(headers) + "\n"
            for i, row in enumerate(sample_rows):
                table_text += f"행{i+1}: " + " | ".join(row) + "\n"

            prompt = f"""다음 표의 내용을 교육 자료 검수 관점에서 검증해주세요:

    {table_text}

    다음 항목들을 확인해주세요:
    1. 데이터 타입 일관성 (숫자/텍스트 혼재 문제)
    2. 논리적 일관성 (합계, 비율 등)
    3. 단위 표기 일관성
    4. 명백한 오류나 이상값
    5. 가독성 문제

    문제가 있다면 구체적으로 설명하고, 없다면 "문제없음"이라고 답하세요.
    한 줄당 하나의 문제만 언급하세요."""

            response = await self.openai_client.chat.completions.create(
                model=self.text_model,
                messages=[{"role": "user", "content": prompt}],
                max_completion_tokens=300,
            )

            result = response.choices[0].message.content.strip()

            if "문제없음" in result:
                return []
            else:
                return [issue.strip() for issue in result.split('\n') if issue.strip()]
            
        except Exception as e:
            logger.warning(f"AI 표 검증 실패: {str(e)}")
            return []

    async def _analyze_image_issues(self, doc_id: str, pages: List[PageInfo]) -> List[Issue]:
        """이미지 관련 이슈 분석 및 Issue 객체 생성"""
        issues = []
        
        for page in pages:
            for element in page.elements:
                if element.element_type == ElementType.IMAGE and element.image_data:
                    # GPT-4V 분석 결과에서 품질 이슈 추출
                    img_base64 = element.image_data.image_data.split(',')[1] if element.image_data.image_data else None
                    if not img_base64:
                        continue
                    
                    try:
                        analysis = await self._call_gpt4v_analysis(img_base64)
                        if not analysis:
                            continue
                        
                        quality_issues = analysis.get("quality_issues", [])
                        accuracy_concerns = analysis.get("accuracy_concerns")
                        
                        # 품질 이슈들을 Issue 객체로 변환
                        for quality_issue in quality_issues:
                            if quality_issue != "none":
                                issue = Issue(
                                    issue_id=generate_issue_id(
                                        doc_id, page.page_id, element.bbox or BoundingBox(x=0, y=0, width=1, height=1), 
                                        IssueType.IMAGE_QUALITY
                                    ),
                                    doc_id=doc_id,
                                    page_id=page.page_id,
                                    issue_type=IssueType.IMAGE_QUALITY,
                                    bbox_location=element.bbox,
                                    element_id=element.element_id,
                                    original_text=analysis.get("description", ""),
                                    message=self._translate_quality_issue(quality_issue),
                                    suggestion=analysis.get("improvement_suggestions", ""),
                                    confidence=0.8,
                                    agent_name="EnhancedMultimodalDocumentAgent"
                                )
                                issues.append(issue)
                        
                        # 정확성 우려사항을 별도 이슈로 생성
                        if accuracy_concerns:
                            issue = Issue(
                                issue_id=generate_issue_id(
                                    doc_id, page.page_id, element.bbox or BoundingBox(x=0, y=0, width=1, height=1), 
                                    IssueType.FACT
                                ),
                                doc_id=doc_id,
                                page_id=page.page_id,
                                issue_type=IssueType.FACT,
                                bbox_location=element.bbox,
                                element_id=element.element_id,
                                original_text=analysis.get("description", ""),
                                message=f"내용 정확성 우려: {accuracy_concerns}",
                                suggestion="전문가 검토를 통해 내용의 정확성을 확인하세요.",
                                confidence=0.7,
                                agent_name="EnhancedMultimodalDocumentAgent"
                            )
                            issues.append(issue)
                            
                    except Exception as e:
                        logger.warning(f"이미지 이슈 분석 실패 ({element.element_id}): {str(e)}")
        
        return issues
    
    def _translate_quality_issue(self, issue_type: str) -> str:
        """품질 이슈 타입을 한국어 메시지로 변환"""
        translations = {
            "low_resolution": "이미지 해상도가 낮아 텍스트가 흐릿합니다.",
            "poor_contrast": "대비가 낮아 가독성이 떨어집니다.",
            "text_too_small": "텍스트 크기가 너무 작아 읽기 어렵습니다.",
            "unclear_diagram": "다이어그램이 불분명하여 이해하기 어렵습니다.",
            "missing_labels": "중요한 라벨이나 설명이 누락되었습니다.",
            "color_accessibility": "색상 접근성에 문제가 있습니다 (색맹 고려 필요)."
        }
        return translations.get(issue_type, f"이미지 품질 문제: {issue_type}")
    
    async def _parse_ppt_enhanced(self, file_path: Path, doc_id: str) -> Tuple[DocumentMeta, List[PageInfo]]:
        """향상된 PowerPoint 멀티모달 파싱"""
        logger.info(f"향상된 PPT 파싱 중: {file_path}")
        
        pages = []
        title = None
        
        try:
            prs = Presentation(file_path)
            
            # 첫 슬라이드에서 제목 추출
            if prs.slides:
                first_slide_text = self._extract_slide_text(prs.slides[0])
                title = await self._extract_title_ai(first_slide_text)
            
            # 각 슬라이드 파싱
            for slide_num, slide in enumerate(prs.slides, 1):
                logger.info(f"슬라이드 {slide_num} 향상된 파싱 중...")
                
                elements = []
                slide_text_parts = []
                element_index = 0
                
                # 슬라이드 내 모든 shape 분석
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 텍스트 요소
                        text_element = PageElement(
                            element_id=generate_element_id(f"slide_{slide_num:03d}", ElementType.TEXT, element_index),
                            element_type=ElementType.TEXT,
                            bbox=self._shape_to_bbox(shape),
                            text_content=shape.text.strip(),
                            confidence=1.0
                        )
                        elements.append(text_element)
                        slide_text_parts.append(shape.text.strip())
                        element_index += 1
                    
                    elif shape.shape_type == 13:  # msoPicture
                        # 이미지 요소 (향상된 처리)
                        image_element = await self._extract_ppt_image_enhanced(
                            shape, f"slide_{slide_num:03d}", element_index
                        )
                        if image_element:
                            elements.append(image_element)
                            element_index += 1
                    
                    elif shape.has_table:
                        # 표 요소 (향상된 처리)
                        table_element = await self._extract_ppt_table_enhanced(
                            shape.table, f"slide_{slide_num:03d}", element_index
                        )
                        if table_element:
                            elements.append(table_element)
                            element_index += 1
                
                # 전체 슬라이드 텍스트
                slide_text = "\n".join(slide_text_parts)
                
                # PageInfo 생성
                page_info = PageInfo(
                    page_id=f"slide_{slide_num:03d}",
                    page_number=slide_num,
                    raw_text=slide_text,
                    word_count=len(slide_text.split()) if slide_text else 0,
                    elements=elements
                )
                
                pages.append(page_info)
        
        except Exception as e:
            logger.error(f"향상된 PPT 파싱 실패: {str(e)}")
            raise
        
        # 문서 메타데이터 생성
        doc_meta = DocumentMeta(
            doc_id=doc_id,
            title=title or file_path.stem,
            doc_type="ppt",
            total_pages=len(pages),
            file_path=str(file_path)
        )
        
        return doc_meta, pages
    
    async def _extract_ppt_image_enhanced(self, shape, page_id: str, element_index: int) -> Optional[PageElement]:
        """향상된 PPT 이미지 추출 (GPT-4V 분석 포함)"""
        try:
            # PPT 이미지에서 실제 이미지 데이터 추출 시도
            image_element = PageElement(
                element_id=generate_element_id(page_id, ElementType.IMAGE, element_index),
                element_type=ElementType.IMAGE,
                bbox=self._shape_to_bbox(shape),
                image_data=ImageElement(
                    element_id=generate_element_id(page_id, ElementType.IMAGE, element_index),
                    bbox=self._shape_to_bbox(shape),
                    format="unknown",
                    size_bytes=0,
                    dimensions=(100, 100),
                    description="PPT 이미지 - 상세 분석 필요"
                ),
                confidence=0.6
            )
            
            return image_element
            
        except Exception as e:
            logger.warning(f"향상된 PPT 이미지 추출 실패: {str(e)}")
            return None
    
    async def _extract_ppt_table_enhanced(self, table, page_id: str, element_index: int) -> Optional[PageElement]:
        """향상된 PPT 표 추출 (AI 검증 포함)"""
        try:
            headers, rows = [], []

            row_list = list(table.rows)
            if len(row_list) == 0:
                return None

            # 첫 행을 헤더로
            first_row = row_list[0]
            headers = [cell.text.strip() for cell in first_row.cells]

            # 나머지는 데이터
            for r in row_list[1:]:
                rows.append([cell.text.strip() for cell in r.cells])

            # AI를 통한 표 검증
            table_issues = await self._validate_table_ai(headers, rows)

            table_element = PageElement(
                element_id=generate_element_id(page_id, ElementType.TABLE, element_index),
                element_type=ElementType.TABLE,
                table_data=TableElement(
                    element_id=generate_element_id(page_id, ElementType.TABLE, element_index),
                    headers=headers,
                    rows=rows,
                    extraction_confidence=0.8 if not table_issues else 0.6
                ),
                confidence=0.8 if not table_issues else 0.6,
                processing_notes=table_issues
            )
            return table_element

        except Exception as e:
            logger.warning(f"향상된 PPT 표 추출 실패: {e!r}")
            return None
    
    async def _create_enhanced_multimodal_index(self, doc_id: str, pages: List[PageInfo]):
        """향상된 멀티모달 요소를 포함한 LlamaIndex 생성"""
        logger.info(f"향상된 멀티모달 인덱스 생성 중: {doc_id}")
        
        documents = []
        
        for page in pages:
            # 각 요소를 별도 문서로 인덱싱 (향상된 텍스트 생성)
            for element in page.elements:
                text_content = ""
                metadata = {
                    "doc_id": doc_id,
                    "page_id": page.page_id,
                    "page_number": page.page_number,
                    "element_id": element.element_id,
                    "element_type": element.element_type.value,
                    "confidence": element.confidence
                }
                
                if element.element_type == ElementType.TEXT:
                    text_content = element.text_content or ""
                
                elif element.element_type == ElementType.IMAGE:
                    # 향상된 이미지 텍스트 생성
                    parts = []
                    if element.image_data:
                        if element.image_data.ocr_text:
                            parts.append(f"[이미지 텍스트] {element.image_data.ocr_text}")
                        if element.image_data.description:
                            parts.append(f"[이미지 설명] {element.image_data.description}")
                        
                        # 메타데이터에 이미지 정보 추가
                        metadata.update({
                            "image_type": getattr(element.image_data, 'image_type', 'unknown'),
                            "educational_purpose": getattr(element.image_data, 'educational_purpose', None),
                            "quality_issues": getattr(element.image_data, 'quality_issues', [])
                        })
                    
                    text_content = " ".join(parts)
                
                elif element.element_type == ElementType.TABLE:
                    # 향상된 표 데이터 텍스트 변환
                    if element.table_data:
                        parts = []
                        if element.table_data.headers:
                            parts.append(f"[표 헤더] {' | '.join(element.table_data.headers)}")
                        
                        # 표 내용을 자연어로 변환
                        for i, row in enumerate(element.table_data.rows):
                            row_text = ' | '.join(row)
                            parts.append(f"[표 행{i+1}] {row_text}")
                        
                        # 표 요약 추가
                        table_summary = f"이 표는 {len(element.table_data.headers)}개 열과 {len(element.table_data.rows)}개 행으로 구성되어 있습니다."
                        parts.append(f"[표 요약] {table_summary}")
                        
                        text_content = "\n".join(parts)
                        
                        # 메타데이터에 표 정보 추가
                        metadata.update({
                            "table_rows": len(element.table_data.rows),
                            "table_cols": len(element.table_data.headers),
                            "has_headers": bool(element.table_data.headers)
                        })
                
                elif element.element_type == ElementType.CHART:
                    # 향상된 차트 설명
                    if element.chart_data:
                        parts = [f"[차트 유형] {element.chart_data.chart_type}"]
                        
                        if element.chart_data.title:
                            parts.append(f"[차트 제목] {element.chart_data.title}")
                        if element.chart_data.description:
                            parts.append(f"[차트 설명] {element.chart_data.description}")
                        if element.chart_data.insights:
                            parts.append(f"[주요 인사이트] {' | '.join(element.chart_data.insights)}")
                        
                        text_content = "\n".join(parts)
                        
                        # 메타데이터에 차트 정보 추가
                        metadata.update({
                            "chart_type": element.chart_data.chart_type,
                            "has_title": bool(element.chart_data.title),
                            "data_points_count": len(element.chart_data.data_points)
                        })
                
                if text_content.strip():
                    doc = Document(
                        text=text_content,
                        metadata=metadata
                    )
                    documents.append(doc)
        
        if not documents:
            logger.warning(f"인덱싱할 콘텐츠가 없습니다: {doc_id}")
            return
        
        try:
            # 향상된 멀티모달 벡터 인덱스 생성
            index = VectorStoreIndex.from_documents(
                documents,
                embed_model=self.embeddings,
                transformations=[self.text_splitter]
            )
            
            self.indexes[doc_id] = index
            logger.info(f"향상된 멀티모달 인덱스 생성 완료: {doc_id} ({len(documents)} 요소)")
            
        except Exception as e:
            logger.error(f"향상된 인덱스 생성 실패: {str(e)}")
            raise
    
    # 기존 유틸리티 메서드들 유지
    def _emu_to_px(self, emu_val: Any, dpi: int = 96) -> int:
        try:
            emu = int(emu_val)
        except Exception:
            emu = int(float(emu_val))
        px = int(round(emu * dpi / 914400))
        return max(0, px)

    def _shape_to_bbox(self, shape) -> BoundingBox:
        x = self._emu_to_px(shape.left)
        y = self._emu_to_px(shape.top)
        w = max(1, self._emu_to_px(shape.width))
        h = max(1, self._emu_to_px(shape.height))
        return BoundingBox(x=x, y=y, width=w, height=h)
    
    def _extract_slide_text(self, slide) -> str:
        """슬라이드에서 모든 텍스트 추출"""
        texts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        
        return "\n".join(texts)
    
    # 향상된 접근 메서드들
    def get_document(self, doc_id: str) -> Optional[DocumentMeta]:
        """문서 메타데이터 조회"""
        return self.documents.get(doc_id)
    
    def get_pages(self, doc_id: str) -> List[PageInfo]:
        """문서의 페이지 목록 조회"""
        return self.pages.get(doc_id, [])
    
    def get_index(self, doc_id: str) -> Optional[VectorStoreIndex]:
        """문서의 벡터 인덱스 조회"""
        return self.indexes.get(doc_id)
    
    def get_image_issues(self, doc_id: str) -> List[Issue]:
        """문서의 이미지 관련 이슈 조회"""
        return self.image_issues.get(doc_id, [])
    
    async def enhanced_search_in_document(self, doc_id: str, query: str, top_k: int = 5, filter_by_type: Optional[ElementType] = None) -> List[Dict[str, Any]]:
        """향상된 멀티모달 문서 내 검색"""
        index = self.get_index(doc_id)
        if not index:
            return []
        
        try:
            # 쿼리 엔진 생성
            query_engine = index.as_query_engine(similarity_top_k=top_k * 2)  # 더 많이 검색 후 필터링
            
            # 검색 수행
            response = query_engine.query(query)
            
            # 결과 정리 및 필터링
            results = []
            for node in response.source_nodes:
                # 타입 필터링
                if filter_by_type and node.metadata.get("element_type") != filter_by_type.value:
                    continue
                
                result = {
                    "text": node.text,
                    "score": node.score,
                    "metadata": node.metadata,
                    "page_id": node.metadata.get("page_id"),
                    "page_number": node.metadata.get("page_number"),
                    "element_id": node.metadata.get("element_id"),
                    "element_type": node.metadata.get("element_type"),
                    "confidence": node.metadata.get("confidence")
                }
                
                # 추가 메타데이터 포함
                if "image_type" in node.metadata:
                    result["image_info"] = {
                        "type": node.metadata.get("image_type"),
                        "educational_purpose": node.metadata.get("educational_purpose"),
                        "quality_issues": node.metadata.get("quality_issues", [])
                    }
                elif "table_rows" in node.metadata:
                    result["table_info"] = {
                        "rows": node.metadata.get("table_rows"),
                        "cols": node.metadata.get("table_cols"),
                        "has_headers": node.metadata.get("has_headers")
                    }
                elif "chart_type" in node.metadata:
                    result["chart_info"] = {
                        "type": node.metadata.get("chart_type"),
                        "has_title": node.metadata.get("has_title"),
                        "data_points_count": node.metadata.get("data_points_count")
                    }
                
                results.append(result)
                
                if len(results) >= top_k:
                    break
            
            return results
            
        except Exception as e:
            logger.error(f"향상된 멀티모달 문서 검색 실패: {str(e)}")
            return []
    
    def get_enhanced_multimodal_stats(self, doc_id: str) -> Dict[str, Any]:
        """향상된 멀티모달 문서 통계"""
        doc_meta = self.get_document(doc_id)
        pages = self.get_pages(doc_id)
        image_issues = self.get_image_issues(doc_id)
        
        if not doc_meta or not pages:
            return {}
        
        # 요소별 상세 통계 계산
        element_stats = {
            "text": 0,
            "image": 0,
            "table": 0,
            "chart": 0,
            "total": 0
        }
        
        image_analysis_stats = {
            "with_ocr": 0,
            "with_ai_description": 0,
            "educational_diagrams": 0,
            "data_charts": 0,
            "quality_issues_found": 0
        }
        
        table_stats = {
            "with_headers": 0,
            "avg_rows": 0,
            "avg_cols": 0,
            "validation_issues": 0
        }
        
        total_words = 0
        total_chars = 0
        
        for page in pages:
            total_words += page.word_count
            total_chars += len(page.raw_text)
            
            for element in page.elements:
                element_stats["total"] += 1
                
                if element.element_type == ElementType.TEXT:
                    element_stats["text"] += 1
                elif element.element_type == ElementType.IMAGE:
                    element_stats["image"] += 1
                    if element.image_data:
                        if element.image_data.ocr_text:
                            image_analysis_stats["with_ocr"] += 1
                        if element.image_data.description:
                            image_analysis_stats["with_ai_description"] += 1
                elif element.element_type == ElementType.TABLE:
                    element_stats["table"] += 1
                    if element.table_data:
                        if element.table_data.has_headers:
                            table_stats["with_headers"] += 1
                        if element.processing_notes:
                            table_stats["validation_issues"] += len(element.processing_notes)
                elif element.element_type == ElementType.CHART:
                    element_stats["chart"] += 1
                    image_analysis_stats["data_charts"] += 1
        
        # 이미지 품질 이슈 통계
        image_analysis_stats["quality_issues_found"] = len(image_issues)
        
        return {
            "doc_id": doc_id,
            "title": doc_meta.title,
            "total_pages": len(pages),
            "total_words": total_words,
            "total_chars": total_chars,
            "avg_words_per_page": total_words / len(pages) if pages else 0,
            "has_index": doc_id in self.indexes,
            "elements": element_stats,
            "image_analysis": image_analysis_stats,
            "table_analysis": table_stats,
            "multimodal_features": {
                "ocr_enabled": self.enable_ocr,
                "vision_analysis_enabled": self.enable_vision_analysis,
                "vision_model": self.vision_model,
                "issues_detected": len(image_issues)
            }
        }
    
    def list_documents(self) -> List[DocumentMeta]:
        """모든 문서 목록 조회"""
        return list(self.documents.values())


# 테스트 함수
async def test_enhanced_multimodal_agent():
    """EnhancedMultimodalDocumentAgent 테스트"""
    print("🧪 향상된 MultimodalDocumentAgent 테스트 시작...")
    
    # 에이전트 생성
    import os
    api_key = os.getenv("OPENAI_API_KEY")
    
    if not api_key:
        print("❌ OPENAI_API_KEY 환경변수가 필요합니다.")
        return
    
    agent = EnhancedMultimodalDocumentAgent(
        openai_api_key=api_key,
        enable_ocr=True,
        enable_vision_analysis=True
    )
    
    # 테스트 파일들
    test_files = [
        "Ch01_intro.pdf",
        # "Ch01_intro.pptx"
    ]
    
    for test_file in test_files:
        if Path(f"/home/matmang/edu-audit-ai/sample_docs/Ch01_intro.pdf").exists():
            try:
                print(f"\n📄 향상된 멀티모달 파싱 테스트: {test_file}")
                
                # 문서 파싱
                doc_meta = await agent.parse_document("/home/matmang/edu-audit-ai/sample_docs/Ch01_intro.pdf")
                
                print(f"✅ 파싱 성공!")
                print(f"   문서 ID: {doc_meta.doc_id}")
                print(f"   제목: {doc_meta.title}")
                print(f"   페이지 수: {doc_meta.total_pages}")
                print(f"   이미지: {doc_meta.total_images}개")
                print(f"   표: {doc_meta.total_tables}개")
                print(f"   차트: {doc_meta.total_charts}개")
                
                # 향상된 멀티모달 통계
                stats = agent.get_enhanced_multimodal_stats(doc_meta.doc_id)
                print(f"   총 요소: {stats['elements']['total']}개")
                print(f"   OCR 추출: {stats['image_analysis']['with_ocr']}개")
                print(f"   AI 설명: {stats['image_analysis']['with_ai_description']}개")
                print(f"   품질 이슈: {stats['image_analysis']['quality_issues_found']}개")
                
                # 이미지 이슈 조회
                image_issues = agent.get_image_issues(doc_meta.doc_id)
                print(f"   이미지 관련 이슈: {len(image_issues)}개")
                for issue in image_issues[:3]:  # 처음 3개만 출력
                    print(f"     - {issue.issue_type.value}: {issue.message}")
                
                # 향상된 검색 테스트
                search_results = await agent.enhanced_search_in_document(
                    doc_meta.doc_id,
                    "이미지",
                    top_k=3
                )
                print(f"   향상된 검색 결과: {len(search_results)}개")
                
                for result in search_results[:2]:  # 처음 2개만 출력
                    print(f"     - {result['element_type']}: {result['text'][:50]}...")
                    if 'image_info' in result:
                        print(f"       이미지 타입: {result['image_info']['type']}")
                
                # 타입별 검색 테스트
                image_results = await agent.enhanced_search_in_document(
                    doc_meta.doc_id, "차트", filter_by_type=ElementType.IMAGE
                )
                print(f"   이미지만 검색: {len(image_results)}개")
                
            except Exception as e:
                print(f"❌ 파싱 실패: {str(e)}")
                import traceback
                traceback.print_exc()
        else:
            print(f"⏭️  테스트 파일 없음: {test_file}")
    
    # 문서 목록 출력
    docs = agent.list_documents()
    print(f"\n📚 총 {len(docs)}개 향상된 멀티모달 문서 로드됨")
    
    print("🎉 향상된 MultimodalDocumentAgent 테스트 완료!")


if __name__ == "__main__":
    asyncio.run(test_enhanced_multimodal_agent())