"""
EDU-Audit Document Agent - Simplified
슬라이드 이미지 → GPT 캡션 → 임베딩 → 라마인덱스 파이프라인
"""

import asyncio
import logging
import os
import base64
import json
import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Set
from io import BytesIO
from collections import defaultdict

# PDF/PPT 파싱 라이브러리
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image

# OpenAI API
import openai
from openai import AsyncOpenAI

# LlamaIndex 관련
from llama_index.core import Document, VectorStoreIndex
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI as LlamaOpenAI

# 기본 모델들만 사용
from src.core.models import DocumentMeta, generate_doc_id

from dotenv import load_dotenv

env_path = Path(__file__).resolve().parents[2] / '.env.dev'
load_dotenv(env_path)

logger = logging.getLogger(__name__)

class DocumentAgent:
    """
    용어/심볼 딕셔너리 기능이 추가된 DocumentAgent
    """
    
    def __init__(
        self, 
        openai_api_key: Optional[str] = None,
        vision_model: str = "gpt-5-nano",
        embedding_model: str = "text-embedding-3-small",
        image_quality: str = "high"
    ):
        self.openai_api_key = openai_api_key or os.getenv("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError("OpenAI API Key가 필요합니다.")
        
        self.vision_model = vision_model
        self.embedding_model = embedding_model
        self.image_quality = image_quality
        
        # OpenAI 클라이언트
        self.client = AsyncOpenAI(api_key=self.openai_api_key)
        
        # LlamaIndex 컴포넌트
        self.embeddings = OpenAIEmbedding(
            model=embedding_model,
            api_key=self.openai_api_key
        )
        
        # 메모리 저장소
        self.documents: Dict[str, DocumentMeta] = {}
        self.indexes: Dict[str, VectorStoreIndex] = {}
        self.slide_images: Dict[str, List[Dict[str, Any]]] = {}
        
        logger.info(f"DocumentAgent 초기화 완료")
        logger.info(f"  Vision Model: {vision_model}")
        logger.info(f"  Embedding Model: {embedding_model}")
    
    async def process_document(self, file_path: str) -> DocumentMeta:
        """
        문서 처리 메인 파이프라인 (용어 딕셔너리 포함)
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")
        
        logger.info(f" 문서 처리 시작: {file_path}")
        
        # 1. 문서 ID 생성
        doc_id = generate_doc_id(str(file_path))
        
        # 2. 슬라이드 이미지 변환
        slide_images = await self._convert_to_slide_images(file_path, doc_id)
        
        # 3. GPT 캡션 생성
        captioned_slides = await self._generate_captions(slide_images)
        
        # 4. Document 객체 생성 및 인덱싱
        documents = await self._create_documents(captioned_slides, doc_id)
        
        # 5. 벡터 인덱스 구축
        index = await self._build_vector_index(documents)
        
        # 6. 문서 메타데이터 생성
        doc_meta = DocumentMeta(
            doc_id=doc_id,
            title=self._extract_title(file_path.stem, captioned_slides),
            doc_type=file_path.suffix.lower().replace('.', ''),
            total_pages=len(slide_images),
            file_path=str(file_path)
        )
        
        # 7. 메모리에 저장
        self.documents[doc_id] = doc_meta
        self.indexes[doc_id] = index
        self.slide_images[doc_id] = captioned_slides
        
        logger.info(f" 문서 처리 완료: {doc_id}")
        logger.info(f"  총 {len(slide_images)} 슬라이드 처리됨")
        
        return doc_meta
    
    # 기존 DocumentAgent 메서드들 (동일)
    async def _convert_to_slide_images(self, file_path: Path, doc_id: str) -> List[Dict[str, Any]]:
        """파일을 슬라이드 단위 이미지로 변환"""
        logger.info(f"슬라이드 이미지 변환 중: {file_path}")
        
        slide_images = []
        
        if file_path.suffix.lower() == '.pdf':
            slide_images = await self._convert_pdf_to_images(file_path, doc_id)
        elif file_path.suffix.lower() in ['.ppt', '.pptx']:
            slide_images = await self._convert_ppt_to_images(file_path, doc_id)
        else:
            raise ValueError(f"지원하지 않는 파일 형식: {file_path.suffix}")
        
        logger.info(f"이미지 변환 완료: {len(slide_images)} 슬라이드")
        return slide_images
    
    async def _convert_pdf_to_images(self, file_path: Path, doc_id: str) -> List[Dict[str, Any]]:
        """PDF를 페이지별 이미지로 변환"""
        slides = []
        
        try:
            pdf_doc = fitz.open(file_path)
            
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                
                # 페이지를 이미지로 변환 (300 DPI)
                mat = fitz.Matrix(300/72, 300/72)
                pix = page.get_pixmap(matrix=mat)
                
                # PIL Image로 변환
                img_data = pix.pil_tobytes(format="PNG")
                pil_image = Image.open(BytesIO(img_data))
                
                # Base64 인코딩
                buffered = BytesIO()
                pil_image.save(buffered, format="PNG")
                img_base64 = base64.b64encode(buffered.getvalue()).decode()
                
                slide_data = {
                    "doc_id": doc_id,
                    "page_id": f"p{page_num + 1:03d}",
                    "page_number": page_num + 1,
                    "image_base64": img_base64,
                    "dimensions": pil_image.size,
                    "size_bytes": len(img_data),
                    "caption": None
                }
                
                slides.append(slide_data)
                pix = None
                
        except Exception as e:
            logger.error(f"PDF 변환 실패: {str(e)}")
            raise
        finally:
            if 'pdf_doc' in locals():
                pdf_doc.close()
        
        return slides
    
    async def _convert_ppt_to_images(self, file_path: Path, doc_id: str) -> List[Dict[str, Any]]:
        """PowerPoint를 슬라이드별 이미지로 변환"""
        slides = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # 슬라이드 텍스트 추출
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                # 임시 이미지 생성
                dummy_image = Image.new('RGB', (1920, 1080), color='white')
                
                buffered = BytesIO()
                dummy_image.save(buffered, format="PNG")
                img_base64 = base64.b64encode(buffered.getvalue()).decode()
                
                slide_data = {
                    "doc_id": doc_id,
                    "page_id": f"slide_{slide_num:03d}",
                    "page_number": slide_num,
                    "image_base64": img_base64,
                    "dimensions": (1920, 1080),
                    "size_bytes": len(buffered.getvalue()),
                    "slide_text": slide_text,
                    "caption": None
                }
                
                slides.append(slide_data)
        
        except Exception as e:
            logger.error(f"PPT 변환 실패: {str(e)}")
            raise
        
        return slides
    
    async def _generate_captions(self, slide_images: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """슬라이드 이미지들에 대해 GPT 캡션 생성"""
        logger.info(f"GPT 캡션 생성 중: {len(slide_images)} 슬라이드")
        
        captioned_slides = []
        
        for i, slide_data in enumerate(slide_images):
            try:
                logger.info(f"캡션 생성: {slide_data['page_id']} ({i+1}/{len(slide_images)})")
                
                caption = await self._generate_single_caption(slide_data)
                
                slide_data_with_caption = slide_data.copy()
                slide_data_with_caption["caption"] = caption
                
                captioned_slides.append(slide_data_with_caption)
                
                await asyncio.sleep(0.1)
                
            except Exception as e:
                logger.warning(f"캡션 생성 실패 {slide_data['page_id']}: {str(e)}")
                slide_data_with_caption = slide_data.copy()
                slide_data_with_caption["caption"] = f"슬라이드 {slide_data['page_number']}"
                captioned_slides.append(slide_data_with_caption)
        
        logger.info(f"캡션 생성 완료: {len(captioned_slides)} 슬라이드")
        return captioned_slides
    
    async def _generate_single_caption(self, slide_data: Dict[str, Any]) -> str:
        """단일 슬라이드에 대한 GPT 캡션 생성"""
        
        context = ""
        if "slide_text" in slide_data and slide_data["slide_text"].strip():
            context = f"\n\n슬라이드 텍스트:\n{slide_data['slide_text']}"
        
        prompt = f"""이 교육자료 슬라이드를 2~3문장으로 요약하고 설명하세요.

다음 내용을 포함해주세요:
- 슬라이드의 주요 주제/내용
- 핵심 개념이나 정보
- 교육적 목적이나 의도

간단하고 명확하게 작성해주세요.{context}"""

        try:
            response = await self.client.chat.completions.create(
                model=self.vision_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": prompt
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{slide_data['image_base64']}",
                                    "detail": self.image_quality
                                }
                            }
                        ]
                    }
                ],
            )
            
            caption = response.choices[0].message.content.strip()

            return caption
            
        except Exception as e:
            logger.error(f"GPT 캡션 생성 실패: {str(e)}")
            return f"슬라이드 {slide_data['page_number']} - 교육자료 내용"
    
    async def _create_documents(self, captioned_slides: List[Dict[str, Any]], doc_id: str) -> List[Document]:
        """캡션이 있는 슬라이드들을 LlamaIndex Document 객체로 변환"""
        logger.info(f"Document 객체 생성 중: {len(captioned_slides)} 슬라이드")
        
        documents = []
        
        for slide_data in captioned_slides:
            text_content = slide_data["caption"]
            
            if "slide_text" in slide_data and slide_data["slide_text"].strip():
                text_content += f"\n\n원본 텍스트:\n{slide_data['slide_text']}"
            
            metadata = {
                "doc_id": doc_id,
                "page_id": slide_data["page_id"],
                "page_number": slide_data["page_number"],
                "image_path": None,
                "dimensions": slide_data["dimensions"],
                "size_bytes": slide_data["size_bytes"],
                "has_caption": bool(slide_data["caption"])
            }
            
            doc = Document(
                text=text_content,
                metadata=metadata
            )
            
            documents.append(doc)
        
        logger.info(f"Document 객체 생성 완료: {len(documents)}개")
        return documents
    
    async def _build_vector_index(self, documents: List[Document]) -> VectorStoreIndex:
        """Document 리스트로부터 벡터 인덱스 구축"""
        logger.info(f"벡터 인덱스 구축 중: {len(documents)} 문서")
        
        try:
            index = VectorStoreIndex.from_documents(
                documents,
                embed_model=self.embeddings
            )
            
            logger.info("벡터 인덱스 구축 완료")
            return index
            
        except Exception as e:
            logger.error(f"벡터 인덱스 구축 실패: {str(e)}")
            raise
    
    def _extract_title(self, filename: str, slides: List[Dict[str, Any]]) -> str:
        """문서 제목 추출"""
        if slides and slides[0].get("caption"):
            first_caption = slides[0]["caption"]
            sentences = first_caption.split('.')
            if sentences and len(sentences[0]) < 100:
                return sentences[0].strip()
        
        return filename.replace('_', ' ').replace('-', ' ')
    
    # 기존 조회 메서드들 + 용어 딕셔너리 관련 메서드 추가
    
    def get_document(self, doc_id: str) -> Optional[DocumentMeta]:
        """문서 메타데이터 조회"""
        return self.documents.get(doc_id)
    
    def get_index(self, doc_id: str) -> Optional[VectorStoreIndex]:
        """문서의 벡터 인덱스 조회"""
        return self.indexes.get(doc_id)
    
    def get_slide_data(self, doc_id: str) -> List[Dict[str, Any]]:
        """슬라이드 원본 데이터 조회"""
        return self.slide_images.get(doc_id, [])
    
    def search_in_document(self, doc_id: str, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """문서 내 의미적 검색"""
        index = self.get_index(doc_id)
        if not index:
            return []
        
        try:
            query_engine = index.as_query_engine(similarity_top_k=top_k)
            response = query_engine.query(query)
            
            results = []
            for node in response.source_nodes:
                result = {
                    "text": node.text,
                    "score": getattr(node, 'score', 0.0),
                    "metadata": node.metadata,
                    "page_id": node.metadata.get("page_id"),
                    "page_number": node.metadata.get("page_number")
                }
                results.append(result)
            
            return results
            
        except Exception as e:
            logger.error(f"문서 검색 실패: {str(e)}")
            return []
    
    def get_slide_by_page_id(self, doc_id: str, page_id: str) -> Optional[Dict[str, Any]]:
        """특정 슬라이드 데이터 조회"""
        slides = self.get_slide_data(doc_id)
        
        for slide in slides:
            if slide["page_id"] == page_id:
                return slide
        
        return None
    
    def list_documents(self) -> List[DocumentMeta]:
        """모든 문서 목록 조회"""
        return list(self.documents.values())
    
    def get_document_stats(self, doc_id: str) -> Dict[str, Any]:
        """문서 통계 정보"""
        doc_meta = self.get_document(doc_id)
        slides = self.get_slide_data(doc_id)
        
        if not doc_meta or not slides:
            return {}
        
        # 캡션 통계
        captioned_slides = sum(1 for slide in slides if slide.get("caption"))
        total_caption_chars = sum(len(slide.get("caption", "")) for slide in slides)
        
        # 이미지 통계
        total_image_size = sum(slide.get("size_bytes", 0) for slide in slides)
        avg_dimensions = (
            sum(slide.get("dimensions", [0, 0])[0] for slide in slides) / len(slides),
            sum(slide.get("dimensions", [0, 0])[1] for slide in slides) / len(slides)
        ) if slides else (0, 0)
        
        return {
            "doc_id": doc_id,
            "title": doc_meta.title,
            "doc_type": doc_meta.doc_type,
            "total_slides": len(slides),
            "captioned_slides": captioned_slides,
            "caption_coverage": captioned_slides / len(slides) if slides else 0,
            "total_caption_chars": total_caption_chars,
            "avg_caption_length": total_caption_chars / captioned_slides if captioned_slides else 0,
            "total_image_size_mb": total_image_size / (1024 * 1024),
            "avg_dimensions": avg_dimensions,
            "has_index": doc_id in self.indexes,
        }


# 테스트 함수
async def test_document_agent():
    """DocumentAgent 테스트"""
    print("🧪  DocumentAgent 테스트 시작...")
    
    # 환경변수에서 API 키 확인
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        print("❌ OPENAI_API_KEY 환경변수가 설정되지 않았습니다.")
        return
    
    # 테스트할 파일 찾기
    test_files = [
        "sample_docs/sample.pdf", 
    ]
    
    test_file = None
    for file_name in test_files:
        if Path(file_name).exists():
            test_file = file_name
            break
    
    if not test_file:
        print("❌ 테스트할 PDF 파일이 없습니다.")
        print(f"   다음 파일 중 하나를 준비해주세요: {', '.join(test_files)}")
        return
    
    print(f"📄 테스트 파일: {test_file}")
    
    try:
        #  DocumentAgent 생성
        agent = DocumentAgent(
            openai_api_key=api_key,
            vision_model="gpt-5-nano",
            embedding_model="text-embedding-3-small"
        )
        
        print(f"📖  문서 처리 중: {test_file}")
        doc_meta = await agent.process_document(test_file)
        
        print(f"✅ 문서 처리 완료!")
        print(f"   문서 ID: {doc_meta.doc_id}")
        print(f"   제목: {doc_meta.title}")
        print(f"   슬라이드 수: {doc_meta.total_pages}")
        
        
        # 문서 통계
        stats = agent.get_document_stats(doc_meta.doc_id)
        print(f"\n📊 문서 통계:")
        print(f"   캡션 생성률: {stats['caption_coverage']:.1%}")
        print(f"   평균 캡션 길이: {stats['avg_caption_length']:.0f}자")
        print(f"   총 이미지 크기: {stats['total_image_size_mb']:.1f}MB")
        
        # 검색 테스트
        if stats['has_index']:
            search_results = agent.search_in_document(
                doc_meta.doc_id,
                "개요",
                top_k=3
            )
            print(f"   검색 결과: {len(search_results)}개")
            
            for result in search_results[:2]:
                print(f"     - 페이지 {result['page_number']}: {result['text'][:50]}...")
        
        # 슬라이드 데이터 확인
        slides = agent.get_slide_data(doc_meta.doc_id)
        if slides:
            first_slide = slides[0]
            print(f"   첫 슬라이드 캡션: {first_slide.get('caption', 'None')[:100]}...")
        
        print("\n🎉  DocumentAgent 테스트 완료!")
        
    except Exception as e:
        print(f"❌ 테스트 실패: {str(e)}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    asyncio.run(test_document_agent())